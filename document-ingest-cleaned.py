import boto3
import botocore
from langchain_community.document_loaders import PyPDFLoader
import json
from opensearchpy import OpenSearch
from opensearchpy import RequestsHttpConnection, OpenSearch, AWSV4SignerAuth
import os
import streamlit as st
from pptx import Presentation
import pandas as pd
import tempfile
import shutil



#Valid File Types
extension_to_type = {
    '.pdf': 'PDF',
    '.html': 'HTML',
    '.htm': 'HTML',
    '.doc': 'Word',
    '.docx': 'Word',
    '.ppt': 'PowerPoint',
    '.pptx': 'PowerPoint',
    '.xls': 'Excel',
    '.xlsx': 'Excel',
    '.txt': 'Text',
    # Add more mappings as necessary
}


config = botocore.config.Config(connect_timeout=300, read_timeout=300)
bedrock = boto3.client('bedrock-runtime' , 'us-east-1', config = config)

#Setup Opensearch connectionand clinet
host = 'HOSTNAME' #use Opensearch Serverless host here
region = 'REGION'# set region of you Opensearch severless collection
service = 'aoss'
credentials = boto3.Session().get_credentials() #Use enviroment credentials
auth = AWSV4SignerAuth(credentials, region, service) 

oss_client = OpenSearch(
    hosts = [{'host': host, 'port': 443}],
    http_auth = auth,
    use_ssl = True,
    verify_certs = True,
    connection_class = RequestsHttpConnection,
    pool_maxsize = 20
)

def get_embeddings(bedrock, text):
    body_text = json.dumps({"inputText": text})
    modelId = 'amazon.titan-embed-text-v1'
    accept = 'application/json'
    contentType='application/json'

    response = bedrock.invoke_model(body=body_text, modelId=modelId, accept=accept, contentType=contentType)
    response_body = json.loads(response.get('body').read())
    embedding = response_body.get('embedding')

    return embedding

def index_doc(client, vectors, content, source, page_number):

    try:
        page = int(page_number)+1
    except:
        page = page_number

    indexDocument={
        'vectors': vectors,
        'content': content,
        'source': source,
        'page': page
        }

    response = client.index(
        index = "INDEXNAME", #Use your index 
        body = indexDocument,
    #    id = '1', commenting out for now
        refresh = False
    )
    return response


def process_pdf(file_name, original_file_name):
    loader = PyPDFLoader(file_name)
    pages = loader.load_and_split()
    source = os.path.basename(original_file_name)

    for page in pages:
        content = page.page_content
        metadata  = page.metadata

        page_number = metadata["page"]

        embeddings = get_embeddings(bedrock, content)

        try:
            response = index_doc(oss_client, embeddings, content, source, page_number)
            st.write("Page: " + str(page_number) + " Done")
        except:
            st.write("Error indexing page: " + str(page_number))


def process_ppt(file_path, original_file_name):
    prs = Presentation(file_path)

    for slide_number, slide in enumerate(prs.slides):
        # Extracting text from each shape that contains text on the slide
        content = ' '.join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())

        if content:  # Ensure there's text to process
            embeddings = get_embeddings(bedrock, content)
            try:
                # Indexing the content along with its embeddings using just the filename
                response = index_doc(oss_client, embeddings, content, original_file_name, slide_number + 1)  # Slide number is 1-indexed
                st.write(f"Slide: {slide_number + 1} Done")
            except Exception as e:
                st.write(f"Error indexing slide: {slide_number + 1} - {str(e)}")


def process_excel(file_name, original_file_name, max_length=20000):
    xls = pd.ExcelFile(file_name)

    for index, sheet_name in enumerate(xls.sheet_names):
        df = pd.read_excel(xls, sheet_name=sheet_name)
        header = ' '.join(df.columns.astype(str))  # Convert header row to string and concatenate
        rows = df.astype(str).apply(lambda x: ' '.join(x.dropna().values), axis=1).tolist()
        
        content_blocks = []
        current_block = header

        for row in rows:
            if len(current_block) + len(row) + 1 < max_length:  # +1 for space
                current_block += " " + row
            else:
                content_blocks.append(current_block)
                current_block = header + " " + row  # Start new block with header

        if current_block:  # Add the last block if it has content
            content_blocks.append(current_block)

        for block in content_blocks:
            if block.strip():
                embeddings = get_embeddings(bedrock, block)
                try:
                    response = index_doc(oss_client, embeddings, block, f"{original_file_name} - {sheet_name}", sheet_name)
                    st.write(f"Processed part of worksheet: {sheet_name}")
                except Exception as e:
                    st.write(f"Error processing part of worksheet: {sheet_name} - {str(e)}")


def determine_file_type(file, extension_to_type):
    # Get the filename from the uploaded file
    filename = file.name
    
    # Extract the file extension
    _, file_extension = os.path.splitext(filename)
    
    # Normalize the file extension to lowercase to ensure case-insensitive matching
    file_extension = file_extension.lower()

    # Return the file type based on the file extension, or 'Unknown' if not found
    return extension_to_type.get(file_extension, 'Unknown')



def save_uploaded_file(uploaded_file):
    _, file_extension = os.path.splitext(uploaded_file.name)
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp:
        tmp.write(uploaded_file.getvalue())
        tmp.flush()
        return tmp.name, uploaded_file.name  # Return both the temporary file path and the original file name


def process_files(file_path, original_file_name, file_type):
    if file_type == 'PDF':
        process_pdf(file_path, original_file_name)
    elif file_type == 'PowerPoint':
        process_ppt(file_path, original_file_name)
    elif file_type == 'Excel':
        process_excel(file_path, original_file_name)


#STREAMLIT


# Let the user upload a file
uploaded_file = st.file_uploader("Upload your file", type=list(extension_to_type.keys()))

if uploaded_file is not None:
    # Display the file type
    file_type = determine_file_type(uploaded_file, extension_to_type)
    st.write(f"The uploaded file is a: {file_type}")



    # Add a button to trigger processing
    if st.button('Process File'):
        # Save the uploaded file to a temporary file
        saved_file_path, original_file_name = save_uploaded_file(uploaded_file)
        
        try:
            # Process the file based on its type
            process_files(saved_file_path, original_file_name, file_type)
            st.success("Processing complete!")
        except Exception as e:
            st.error(f"An error occurred during processing: {str(e)}")
        finally:
            # Optionally, delete the temporary file after processing
            os.remove(saved_file_path)